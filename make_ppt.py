"""
make_ppt.py — 生成 vLLM vs SGLang 分享 PPT
骨架(06-10) + Part1/2 内容填充(06-11)。图表用 assets/ 下实测图。
运行: /home/guoda/python/bin/python3 make_ppt.py → assets/vllm_sglang_share.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x2D, 0x6A, 0x3E)
GRAY = RGBColor(0x55, 0x55, 0x55)


def title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.runs[0].font.size = Pt(size)
    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = BLUE


def body(slide, lines, x=0.7, y=1.5, w=7.0, size=18):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run(); run.text = line if line else " "
        run.font.size = Pt(size)
        if line.startswith("→") or line.startswith("★"):
            run.font.color.rgb = GREEN; run.font.bold = True


def img(slide, path, x=7.8, y=1.5, w=5.2):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new():
        return prs.slides.add_slide(blank)

    A = "assets/"

    # 封面
    s = new(); title(s, "vLLM vs SGLang：推理框架深度对比", 36)
    body(s, ["实测数据 · RTX 3060 12GB · Qwen2.5-1.5B",
             "vLLM 0.22.1 / SGLang 0.5.13", "", "2026-07 技术分享"], y=2.5, w=12, size=22)

    # 目录
    s = new(); title(s, "目录")
    body(s, ["① 为什么需要推理框架", "② vLLM：三层优化(PagedAttention/Continuous Batching/APC)",
             "③ SGLang：缓存粒度革新(RadixAttention)", "④ 对比数据(吞吐/延迟/复用率)",
             "⑤ 选型建议"], w=12, size=22)

    # ===== Part 1 =====
    s = new(); title(s, "① GPU 利用率：朴素推理在浪费什么？")
    body(s, ["decode 是访存瓶颈：每生成 1 token 要把整个模型权重",
             "从显存搬一遍(1.5B≈3GB)。",
             "", "实测单请求 decode: 90 tok/s",
             "batch=10 时: 849 tok/s 总吞吐",
             "→ 单请求时 GPU 算力大量闲置，等串行依赖",
             "→ 推理框架的核心价值: 榨干这部分闲置算力"])

    s = new(); title(s, "② KV Cache：推理的核心瓶颈")
    body(s, ["KV Cache 存历史 token 的 Key/Value，避免重算。",
             "但它吃满显存——KV 越大，能并发的 batch 越小。",
             "", "传统方案: 每请求预留 max_len 连续显存",
             "→ 10 个 50-token 请求，浪费 98.8%!",
             "→ KV 管理决定吞吐: 谁能装下更多并发请求，谁吞吐高"])
    img(s, A + "kvcache_diff.png")

    s = new(); title(s, "③ 传统 KV Cache 的内存浪费")
    body(s, ["问题: 预分配 max_len → 外部碎片 + 内部浪费",
             "", "vLLM 之前: 业界要么预留(浪费)，要么频繁搬迁(慢)",
             "→ PagedAttention 借鉴操作系统虚拟内存:",
             "   分页 + 按需分配 + 间接寻址",
             "→ 浪费从 98.8% 压到 21.9%，同显存多服务 65 倍请求"])

    # ===== Part 2: vLLM =====
    s = new(); title(s, "vLLM ①: PagedAttention — 向操作系统学内存管理")
    body(s, ["KV 切成 16-token 块，Block Table 间接寻址",
             "(逻辑连续 / 物理离散)",
             "", "实测(RTX 3060): 14268 块 × 16 = 228,288 token",
             "并发上限 55.73x",
             "→ free list + 引用计数 + LRU + 惰性失效",
             "→ 本质是操作系统页分配器搬到 KV cache"])
    img(s, A + "pagedattention_arch.png", w=5.0)

    s = new(); title(s, "vLLM ②: Continuous Batching — 消灭等待气泡")
    body(s, ["逐 step 调度: 完成的请求立即移除，新请求随到随插",
             "(不像传统 batching 要等最慢的)",
             "", "实测: 并发 1→5 吞吐涨 3 倍(28.9→88.1 RPS)",
             "峰值 90.7 RPS，拐点并发 8-16",
             "→ token_budget 共享: 长 prefill 不饿死 decode"])

    s = new(); title(s, "vLLM ③: APC — 前缀只算一次")
    body(s, ["块哈希(链式): hash(块i)=f(hash(块i-1), tokens)",
             "命中前缀块直接复用，跳过 prefill",
             "", "实测: WARM TTFT 33ms vs OFF 153ms (降 4.6x)",
             "命中率 98.7%，只省 prefill 不影响 decode",
             "→ 但块粒度: 必须 16-token 对齐才命中"])
    img(s, A + "prefix_ratio_comparison.png", w=5.0)

    # ===== Part 3: SGLang =====
    s = new(); title(s, "SGLang ①: RadixAttention — 缓存粒度降到 token")
    body(s, ["radix 树 + page_size=1: 任意 token 位置可 split",
             "前缀逐 token 精确复用(vLLM 是 16-token 块)",
             "", "实测半块前缀命中: SGLang 94.9% vs vLLM 82.1%",
             "→ 多救回 15 token = 一个 block 的对齐损失",
             "→ 缓存粒度才是 SGLang 的本质优势"])
    img(s, A + "radixattention_arch.png", w=5.0)

    s = new(); title(s, "SGLang ②: 缓存感知调度 + Compressed FSM")
    body(s, ["调度: cache-aware(LPM/DFS_WEIGHT)，主动把同前缀集中",
             "(vLLM 是 FCFS)",
             "Compressed FSM: JSON 确定性段一步跳过 → 吞吐+60%",
             "HiCache: KV 可降级 CPU 二级缓存(放不下不丢)",
             "→ SGLang 心疼每一份算过的 KV"])

    # ===== Part 4: 对比 =====
    s = new(); title(s, "对比 ①: 无共享前缀 → 两框架持平")
    body(s, ["固定输出，无前缀复用，3 轮中位数:",
             "", "并发32: vLLM 1443 vs SGLang 1342 tok/s (差 8%)",
             "TTFT/ITL 也持平，拐点都在并发 8-16",
             "→ 无复用时两框架持平，vLLM 略优 + 生态成熟",
             "→ 社区(H100/ShareGPT)验证同样趋势"])
    img(s, A + "p1_throughput.png")

    s = new(); title(s, "对比 ②: 高复用 → SGLang 2.1x (最关键)")
    body(s, ["前缀复用率 0→100% 梯度实验:",
             "", "0% 复用: RPS 7.31 vs 7.08 (持平)",
             "100% 复用: RPS 8.88 vs 18.95",
             "★ SGLang 吞吐 2.1x !",
             "→ 大量请求共享同一长前缀时，radix 树全局共享",
             "→ 成本可减半"])
    img(s, A + "prefix_ratio_comparison.png")

    s = new(); title(s, "对比 ③: 核心差异 — 命中粒度 1 vs 16 token")
    body(s, ["vLLM 块哈希: 16-token 块，整块对齐才命中",
             "SGLang radix: 1-token，任意长度命中",
             "", "→ 共享前缀不对齐块边界时，vLLM 尾部一整块白算",
             "→ 这就是高复用 SGLang 占优的根源"])
    img(s, A + "kvcache_diff.png")

    # ===== Part 5: 选型 =====
    s = new(); title(s, "选型建议: 复用率决定选型")
    body(s, ["通用/低复用(批量翻译/独立摘要):",
             "  → vLLM (略优 + 生态成熟)",
             "高复用(统一 system prompt/RAG/Agent):",
             "  → SGLang (吞吐 2.1x，成本减半)",
             "延迟/结构化输出/量化: 两者相当",
             "→ 看负载特征，不看'哪个更好'"])
    img(s, A + "selection_guide.png", w=5.2)

    s = new(); title(s, "一句话结论")
    body(s, ["", "KV Cache 管理(块 vs token)是两框架最大差异。",
             "", "★ 复用率是分水岭:",
             "   无复用持平(vLLM略优) → 高复用 SGLang 2.1x",
             "", "通用选 vLLM，高复用选 SGLang。"], w=12, size=22)

    s = new(); title(s, "谢谢 · Q&A")
    body(s, ["实测环境: RTX 3060 12GB / Qwen2.5-1.5B",
             "数据: 3 轮中位数，社区交叉验证",
             "", "局限: 趋势可外推，绝对数字需目标硬件实测"], w=12, size=20)

    out = "assets/vllm_sglang_share.pptx"
    prs.save(out)
    print(f"saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} 张)")


if __name__ == "__main__":
    main()
